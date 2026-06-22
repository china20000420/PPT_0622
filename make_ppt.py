
"""
PPT 终稿 — 零重叠 · 文字精简 · 公式渲染 · 论文原图
"""
from pptx import Presentation; from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor; from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage
import matplotlib; matplotlib.use('Agg'); import matplotlib.pyplot as plt
import os, io

BASE=r"c:\Users\26024\Desktop\做PPT"
FIGS={"bf":os.path.join(BASE,"02_BestFit_Packing","figures"),
      "ic":os.path.join(BASE,"03_InContext_Pretraining","figures"),
      "dd":os.path.join(BASE,"04_Dataset_Decomposition","figures")}
CF='Microsoft YaHei'; EF='Times New Roman'
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
WH=RGBColor(0xFF,0xFF,0xFF); BL=RGBColor(0x14,0x4D,0xD9); BK=RGBColor(0x1A,0x1A,0x1A)
LG=RGBColor(0xE8,0xEE,0xF7); DL=RGBColor(0x0E,0x38,0xA0); GAP=0.22

def Bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb=WH
def T(s,l,t,w,h,txt,sz=12,b=False,fn=CF,c=None,al=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
    p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c or BK
    p.font.bold=b; p.font.name=fn; p.alignment=al; p.space_after=Pt(4); return tf
def M(s,l,t,w,h,runs,al=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=al; p.space_after=Pt(4)
    fst=True
    for r_ in runs:
        tx,sz,bo=r_[0],r_[1],r_[2]; nl=r_[3] if len(r_)>3 else False
        fn=r_[4] if len(r_)>4 else CF; cl=r_[5] if len(r_)>5 else BK
        if nl and not fst: p=tf.add_paragraph(); p.alignment=al; p.space_after=Pt(4)
        run=p.add_run(); run.text=tx; run.font.size=Pt(sz); run.font.color.rgb=cl
        run.font.bold=bo; run.font.name=fn; fst=False
    return tf
def Ln(s,l,t,w,h,fc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    sh.line.fill.background(); return sh
def Img(s,p,l,t,w,h):
    """h is REQUIRED — no auto-sizing to avoid overlaps"""
    if not os.path.exists(p): return
    s.shapes.add_picture(p,Inches(l),Inches(t),Inches(w),Inches(h))
def Hd(s,ti,su=None):
    Ln(s,0.45,0.38,0.04,0.45,fc=BL);T(s,0.75,0.3,11.5,0.5,ti,sz=26,b=True)
    if su:T(s,0.75,0.7,11.5,0.22,su,sz=10,b=False)
def Pn(s,n):T(s,12.5,7.12,0.55,0.2,str(n),sz=10,b=False,al=PP_ALIGN.RIGHT)
def Bb(s):Ln(s,0,7.33,13.333,0.03,fc=BL)
def Sep(s,y):Ln(s,0.75,y,11.8,0.004,fc=LG)

def Tbl(s,l,t,w,h,data):
    R,C=len(data),len(data[0])
    tbl=s.shapes.add_table(R,C,Inches(l),Inches(t),Inches(w),Inches(h)).table
    for r in range(R):
        for c in range(C):
            cl=tbl.cell(r,c);cl.text=data[r][c]
            for p in cl.text_frame.paragraphs:
                p.font.size=Pt(10);p.font.name=CF
                if r==0:p.font.bold=True;p.font.color.rgb=WH;p.alignment=PP_ALIGN.CENTER
                elif c==0:p.font.bold=True;p.font.color.rgb=BK
                else:p.font.color.rgb=BK
            cl.fill.solid();cl.fill.fore_color.rgb=BL if r==0 else (LG if r%2==0 else WH)

def F(s,latex,l,t,w,fs=13):
    """Returns (bottom_y). Caller MUST add GAP before next element."""
    fig,ax=plt.subplots(figsize=(w*1.2,0.55))
    ax.text(0.5,0.5,f'${latex}$',transform=ax.transAxes,fontsize=fs,ha='center',va='center')
    ax.axis('off')
    buf=io.BytesIO()
    fig.savefig(buf,dpi=220,bbox_inches='tight',pad_inches=0.08,transparent=True,format='png')
    plt.close(fig);buf.seek(0)
    im=PILImage.open(buf);iw_px,ih_px=im.size;h_in=w*ih_px/iw_px
    tmp=os.path.join(BASE,'_tf.png');im.save(tmp)
    s.shapes.add_picture(tmp,Inches(l),Inches(t),Inches(w),Inches(h_in))
    return t+h_in

# ═══════ COVER ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s)
Ln(s,0,0,13.333,7.5,fc=BL);Ln(s,10.5,-0.8,5,5,fc=DL);Ln(s,11.5,3.5,2,2,fc=RGBColor(0x24,0x60,0xEE))
T(s,0.9,1.5,11.5,1.0,'大模型预训练数据拼接策略',sz=48,c=WH,b=True)
T(s,0.9,2.6,11.5,0.4,'From Random Concatenation to Intelligent Composition',sz=17,c=RGBColor(0xBB,0xCC,0xF0),fn=EF)
Ln(s,0.9,4.5,2.0,0.025,fc=WH)
T(s,0.9,4.8,11.5,0.35,'Amazon ICML 2024  ·  Meta ICLR 2024  ·  Apple NeurIPS 2024',sz=15,c=RGBColor(0x99,0xB0,0xE0))
T(s,0.9,5.5,11.5,0.3,'2026 年 6 月',sz=14,c=RGBColor(0x99,0xB0,0xE0))

# ═══════ AGENDA ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,2);Hd(s,'目  录')
ag=[('一','为什么要做数据拼接','自回归训练 · GPU 约束 · 文档截断'),
    ('二','随机拼接 (Concat-and-Chunk)','基线做法 · 截断统计 · 三个缺陷'),
    ('三','Best-Fit Packing','Amazon ICML 2024 · 装箱建模 · BFD · 线段树'),
    ('四','In-Context Pretraining','Meta ICLR 2024 · 语义嵌入 · 图遍历'),
    ('五','Dataset Decomposition','Apple NeurIPS 2024 · 二进制分解 · VSL 课程'),
    ('六','对比与总结','方法论·性能对比·演进时间线')]
for i,(n,ti,de) in enumerate(ag):
    y=1.45+i*0.95;T(s,1.0,y+0.02,0.55,0.45,n,sz=26,b=True)
    T(s,1.7,y+0.05,9,0.35,ti,sz=17,b=True);T(s,1.7,y+0.4,10,0.28,de,sz=12,b=False)
    if i<5:Sep(s,y+0.82)

# ═══════ 3. WHY 1 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,3)
Hd(s,'一、为什么要做数据拼接 (1/2)','自回归机制、自注意力计算与 GPU 约束')

# Left — formulas
y=1.3
y=F(s,r'L(\theta)=-\frac{1}{N}\sum_{t=1}^{N}\log P(x_t\mid x_{<t};\theta)',0.75,y,6.0,fs=15)
y+=GAP
y=F(s,r'\mathrm{Attention}(Q,K,V)=\mathrm{softmax}\!\left(\frac{QK^T}{\sqrt{d_k}}\right)V',0.75,y,6.0,fs=14)
y+=GAP
y=F(s,r'Q=XW_Q,\ K=XW_K,\ V=XW_V\in\mathbb{R}^{L\times d}',0.75,y,6.0,fs=13)
y+=GAP
M(s,0.75,y,6.0,2.0,[
    ('QK^T ∈ R^{L×L}。计算与存储均为 O(L²)。L:2048→8192，成本变为 16 倍。',12,True,True),
    ('因果掩码：token t 只关注 1..t−1，但不感知文档边界。',12,True,True),
    ('→ 这是跨文档注意力噪声的直接来源。',12,True,True),
    ('',4,False,True),
    ('GPU 要求 batch 内所有序列 shape 一致。变长文档统一到 L 有两种方式：',12,True,False),
    ('· Padding：填 [PAD]，mask 掉损失但 GPU 仍计算。利用率≈有效 token 占比。',11,False,True),
    ('· Packing (拼接)：多文档拼为等长序列，零填充，利用率 100%。',11,False,True),
    ('  → Decoder-only 架构 (GPT 类) 的标准选择。',11,True,True),
])

# Right
y=1.3
M(s,7.0,y,5.8,5.8,[
    ('RefinedWeb 文档长度分布 (Ding et al., 2024)：',13,True,False),
    ('',3,False,True),
    ('· 平均长度 ≈ 600 tokens',11,False,True),
    ('· 80%+ 文档 < 1000 tokens',11,False,True),
    ('· 约 60% 文档 < 500 tokens',11,False,True),
    ('· 仅约 5% 文档 > 2048 tokens',11,False,True),
    ('→ L=2048 时绝大多数文档远短于 L，必须拼接。',12,True,True),
    ('',6,False,True),
    ('拼接策略直接决定三个核心指标：',13,True,False),
    ('',3,False,True),
    ('① 文档截断率',12,True,True),
    ('   多少文档的完整性在训练中被破坏。',10,False,True),
    ('',3,False,True),
    ('② 跨文档注意力噪声',12,True,True),
    ('   多少 GPU 计算被浪费在无关文档之间。',10,False,True),
    ('',3,False,True),
    ('③ 模型幻觉率',12,True,True),
    ('   Grounding context 丢失导致的生成问题。',10,False,True),
    ('',6,False,True),
    ('本文四种方法在这三个指标上做出了不同的权衡。',13,True,True),
])

# ═══════ 4. TRUNCATION ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,4)
Hd(s,'一、为什么要做数据拼接 (2/2)','Ding et al., ICML 2024, Figure 2')

# Three images at top — fixed heights to avoid overlap
Img(s,os.path.join(FIGS["bf"],"concat_ex1.png"),0.2,1.2,3.95,3.8)
Img(s,os.path.join(FIGS["bf"],"concat_ex2.png"),4.55,1.2,3.95,3.8)
Img(s,os.path.join(FIGS["bf"],"concat_ex3.png"),8.9,1.2,3.95,3.8)
T(s,0.2,5.05,3.95,0.6,'(a) Undefined Names：代码变量定义与使用跨序列分离。自注意力不跨序列→标识符无定义。',sz=9,b=False)
T(s,4.55,5.05,3.95,0.6,'(b) Ungrounded Content：时间锚点与摘要句跨序列。序列 B 无法将生成锚定到可见上下文。',sz=9,b=False)
T(s,8.9,5.05,3.95,0.6,'(c) Missing Knowledge：会议名与举办地跨序列。完整事实无法从任何单一序列学到。',sz=9,b=False)
T(s,0.2,5.7,12.5,0.15,'Figure 2. Ding et al., ICML 2024.',sz=9,b=False,fn=EF)

# Bottom — chart + theory side by side, with gap
Img(s,os.path.join(FIGS["bf"],"truncation_nl_count.png"),0.3,5.95,5.8,1.2)
T(s,0.3,5.85,5.8,0.1,'Figure 4: NL 截断率统计',sz=7,b=False,fn=EF)
M(s,6.5,5.95,6.3,1.2,[
    ('理论分析 (Ding et al., Sec 2.1)：',11,True,False),
    ('文档中位置 m 后 token 依赖前 m token 作为 grounding context。若在 k<m 截断，完整模型 A 与截断模型 B 的期望损失差 ΔL(m)>0，对 m 单调递增。',10,False,True),
    ('→ 截断不是随机扰动，而是对文档尾部产生系统性偏差。',11,True,True),
])

# ═══════ 5. RANDOM 1/3 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,5)
Hd(s,'二、随机拼接 (1/3)','Concat-and-Chunk：GPT-3 / LLaMA / PaLM 等所有主流 LLM 的默认做法')

y=1.3
M(s,0.75,y,5.8,1.0,[('三步操作生成训练序列：',13,True,False),
    ('(1) Fisher-Yates shuffle → 随机排列 π',12,True,True),
    ('(2) 按 π 顺序 tokenize + EOS → 全局 token 流 T',12,True,True),
    ('(3) T 每满 L token 切为一个序列：',12,True,True)])
y=F(s,r'T=\mathrm{tok}(d_{\pi_1})\circ[\mathrm{EOS}]\circ\cdots\circ\mathrm{tok}(d_{\pi_N})\circ[\mathrm{EOS}]',0.75,2.3,5.8,fs=11)
y+=GAP
y=F(s,r's_k=T[kL:(k+1)L],\quad M=\lfloor(\sum l_i+N)/L\rfloor',0.75,y,5.8,fs=14)
y+=GAP
M(s,0.75,y,5.8,2.5,[
    ('每序列严格 L token，零 padding，GPU 利用率 100%。',12,True,True),('',3,False,True),
    ('截断率实测 (RefinedWeb): L=512→85%, 1024→72%, 2048→60%, 4096→45%, 8192→30%',11,False,True,EF),
    ('增大 L 可减截断，但 O(L²) 成本激增。问题未根本解决。',11,True,True),('',3,False,True),
    ('GPT-3 (2020, 175B)→LLaMA (2023, 65B)→PaLM (2022, 540B) 均采用此方案。',10,False,True),
    ('五年间被视为"已解决"的工程细节。直到 Ding et al. (2024) 受控实验系统量化了截断代价。',10,True,True),
])

M(s,7.0,1.3,5.8,5.8,[
    ('三个根本缺陷 (Ding et al., 2024)：',14,True,False),('',4,False,True),
    ('缺陷一：文档完整性系统性破坏',13,True,True),
    ('L=2048 下约 60% 文档以不完整形式参训。代码中变量定义/使用分离→Undefined Name 错误。自然语言因果链切断。大多数训练 token 来自不完整的文档片段。',10,False,True),('',3,False,True),
    ('缺陷二：跨文档注意力噪声',13,True,True),
    ('因果掩码不感知文档边界。序列中第 j 文档的 token 无条件 attend 前 j−1 个随机文档。无预测信号。L 越大噪声越高。',10,False,True),('',3,False,True),
    ('缺陷三：幻觉风险增大',13,True,True),
    ('Grounding context 与结论跨序列分离→模型被训练为"在缺少充分上下文时也做预测"→推理时忽略给定上下文。',10,False,True),('',4,False,True),
    ('以下三种方法分别从算法/语义/范式三层递进回应这些缺陷。',13,True,True),
])

# ═══════ 6. RANDOM 2/3 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,6)
Hd(s,'随机拼接 (2/3)','跨文档注意力噪声的定量分析')

M(s,0.75,1.3,6.2,5.8,[
    ('噪声来源',13,True,False),('',3,False,True),
    ('标准因果掩码：M[t,j] = −∞ if j>t else 0。token t 可以 attend 到所有 j<t 的 token，无论来自哪个文档。掩码对文档边界完全无感知。',11,False,True),('',4,False,True),
    ('定量估计',13,True,False),('',3,False,True),
    ('设长度为 L 的序列包含 K≈L/E[l]≈2048/600≈3.4 个文档。',11,False,True),
    ('对第 k 文档的 token，前 k−1 文档的 token 被无条件 attend。',11,False,True),
    ('噪声占比 ≈ (k−1)/K ≈ 70% (对最后一个文档)。',11,False,True),
    ('对所有 token 平均后，论文保守估计 15–30% 注意力无预测信号。',11,True,True),('',4,False,True),
    ('影响',13,True,False),('',3,False,True),
    ('代码数据：前序可能是不同编程语言→完全无用的注意力。',11,False,True),
    ('知识数据：前序来自不同知识领域→可能产生错误的知识关联。',11,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('缺陷三详解：幻觉风险的因果链',13,True,False),('',4,False,True),
    ('Step 1 — Grounding context 缺失',12,True,True),
    ('"ICML 2024 will be held in Vienna" 跨两个序列。序列 A 仅有 "ICML 2024 will be held"，序列 B 仅有 "in Vienna"。',10,False,True),('',3,False,True),
    ('Step 2 — 训练信号扭曲',12,True,True),
    ('序列 A 中模型需预测 "in"，但正确答案在序列 B 中不可见。模型只能泛化预测 "at" 或 "on"——语法合理但事实错误。',10,False,True),('',3,False,True),
    ('Step 3 — 行为内化',12,True,True),
    ('大量截断样本→模型内化"在不确定时给出语法合理但不保证真实的预测"。推理时即使给了完整上下文，模型仍可能忽略。',10,False,True),('',4,False,True),
    ('三个缺陷的递进关系：',12,True,False),
    ('截断 (根因) → 跨文档噪声 (表现) → 幻觉 (后果)。',11,True,True),
    ('Best-Fit 攻击根因，ICLM 消灭噪声，DD 重新定义范式。',11,True,True),
])

# ═══════ 7. RANDOM 3/3 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,7)
Hd(s,'随机拼接 (3/3) — 总结与引出后续')

M(s,0.75,1.5,11.5,2.5,[
    ('随机拼接的核心矛盾：效率最优，质量有损。',15,True,False),('',3,False,True),
    ('优势：零 padding→GPU 利用率 100%。实现极简。均匀采样公平。无需预处理。',11,False,True),
    ('劣势：截断率 60%→大多数样本不完整。注意力噪声 15–30%→算力浪费。幻觉风险→模型内化猜测行为。',11,False,True),('',6,False,True),
    ('三种改进路线构成递进谱系：',15,True,False),('',4,False,True),
    ('路线一 (算法层)：保持"拼成等长序列"框架，用组合优化最小化截断。',13,True,True),
    ('  → 三、Best-Fit Packing (Amazon, ICML 2024)',13,True,True),('',3,False,True),
    ('路线二 (语义层)：在"拼成等长序列"框架内，改变文档排列顺序消灭跨文档噪声。',13,True,True),
    ('  → 四、ICLM (Meta, ICLR 2024)',13,True,True),('',3,False,True),
    ('路线三 (范式层)：放弃"拼成等长序列"框架，重新定义数据准备方式。',13,True,True),
    ('  → 五、DD (Apple, NeurIPS 2024)',13,True,True),('',4,False,True),
    ('三条路线从浅到深，不互斥，可组合使用。',13,True,True),
])

# ═══════ 8. BF 1/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,8)
Hd(s,'三、Best-Fit Packing (1/6)','Ding, Wang et al. ICML 2024. Amazon AWS. arXiv:2404.10830')
Img(s,os.path.join(FIGS["bf"],"x1.png"),0.2,1.15,8.6,3.9)
T(s,0.2,5.1,8.6,0.18,'Figure 1. Best-Fit Packing 两阶段流程 vs Concatenation。',sz=9,b=False)
M(s,9.15,1.2,3.9,3.5,[
    ('将拼接建模为装箱问题：文档切为 chunk (l≤L)，装最少 bin。',11,True,False),('',2,False,True),
    ('形式化：',11,True,False)])
F(s,r'\min M\ \ \mathrm{s.t.}\ \sum_{c\in s_j}l(c)\leq L,\ \forall j;\ s_1\cup\cdots\cup s_M=C',9.15,2.5,3.9,fs=11)
M(s,9.15,3.2,3.9,3.8,[
    ('NP-Hard (Karp 1972)。BFD 启发式近似。',10,False,True),
    ('Phase 1: l>L 切 chunk (~5% 文档)。Phase 2: BFD 装箱。',10,False,True),
    ('核心结果：阅读 +4.7% · NLI +9.3% · 上下文跟随 +16.8%',10,True,True),
    ('程序合成 +15.0% · 幻觉 −58.3% · 填充 <0.003%',10,True,True),
    ('22 任务零退化。HuggingFace TRL 已集成 (strategy="ffd")。',10,True,True),
])

# ═══════ 9. BF 2/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,9)
Hd(s,'Best-Fit Packing (2/6) — Segmentation 与 BFD Packing')

M(s,0.75,1.3,6.2,5.8,[
    ('Phase 1 — Segmentation (文档预切分)',14,True,False),('',3,False,True),
    ('l(d) ≤ L → 完整保留为 1 chunk。不做任何切分。',12,True,True),
    ('l(d) > L → 切 ⌈l/L⌉ chunk: [L,...,L,r]，r=l mod L。',12,True,True),
    ('仅约 5% 文档需切分 (必要截断)。其余 95%+ 完整保留。',12,True,True),('',4,False,True),
    ('Phase 2 — BFD Packing',14,True,False),('',3,False,True),
    ('(a) 计数排序降序 O(N)。l(c)∈[1,L] 有界→无需比较排序。',12,True,True),
    ('(b) BFD 贪心：依次处理 chunk，找剩余 ≥l(c) 且最小的 bin。',12,True,True),
    ('(c) 线段树加速 O(log L)：BFD 不区分同容量 bin→按容量值组织。',12,True,True),
    ('',3,False,True),
    ('总复杂度 O(N log L)。L=8192→log₂L=13→近似线性。',12,True,True),
    ('1B docs≈3h (Python, 单 CPU)。2B docs OBFD vs FFD: 2.5× 加速。',11,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('装箱示例 L=8：chunk={6,4,4,3,2,2,1,7,5,8}, Σ=42',12,True,False),('',2,False,True),
    ('降序: [8,7,6,5,4,4,3,2,2,1]',10,False,True,EF),('',1,False,True),
    ('装箱过程:',11,True,True),
    ('8→B1[0]; 7→B2[1]; 6→B3[2]; 5→B4[3]',10,False,True,EF),
    ('4→B5[4]; 4→B5[0] ★完美; 3→B4[0] ★Best-Fit',10,False,True,EF),
    ('2→B3[0]; 2→B6[6]; 1→B2[0]',10,False,True,EF),('',1,False,True),
    ('M=6, 利用率 87.5%。OPT=5。BFD 界: M≤(11/9)OPT+4。',10,False,True),('',3,False,True),
    ('大规模装箱效率 (论文 Table 1)：',13,True,False),
    ('RefinedWeb L=2048: +0.0024% | L=8192: +0.00063%',10,False,True,EF),
    ('The Stack L=2048: +0.0028%',10,False,True,EF),('',2,False,True),
    ('→ 填充几乎为零！训练步骤数 ≈ 基线！',12,True,True),
    ('→ 截断率 60%→~5%，训练成本不变！',12,True,True),
    ('→ 四种方法中唯一"纯收益、零代价"方案。',12,True,True),
])

# ═══════ 10. BF 3/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,10)
Hd(s,'Best-Fit Packing (3/6) — 线段树优化详解')

M(s,0.75,1.3,6.2,5.8,[
    ('线段树加速的核心前提',14,True,False),('',3,False,True),
    ('BFD 在 Step (b) 中不区分具体是哪个 bin——只关心"有没有容量 ≥l(c) 的 bin，其中剩余最小的是多少"。两个剩余容量相同的 bin 对 BFD 完全等价。',11,False,True),('',4,False,True),
    ('因此只需为 [1,L] 这 L 个容量值维护数据结构，而非为 N 个 bin。',12,True,True),('',4,False,True),
    ('数据结构 — 线段树 (Segment Tree)：',14,True,False),('',3,False,True),
    ('· L 个叶节点对应容量 1..L。第 i 叶值=i (若存在该容量 bin)，否则=-∞。',11,False,True),
    ('· 内部节点 = max(左子, 右子)。根节点存储所有叶中的最大值。',11,False,True),('',3,False,True),
    ('查询操作 (找 Best-Fit)：从根出发，若左子最大值 ≥l(c)→进左 (左边有容量够的 bin)。否则→进右。到达叶节点时，叶索引即为 Best-Fit 容量。',11,False,True),('',3,False,True),
    ('更新操作：chunk 放入后，旧容量的叶值改为 -∞，新容量的叶值设为新容量。同时更新从叶到根的路径上所有内部节点。',11,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('复杂度对比',14,True,False),('',3,False,True),
    ('朴素 BFD：排序 O(N log N) + 装箱 N×O(M) ≈ O(N²)。不可行。',11,False,True),
    ('标准 BFD：排序 O(N log N) + 装箱 N×O(log N)=O(N log N)。',11,False,True),
    ('优化 BFD：排序 计数 O(N) + 装箱 N×O(log L)=O(N log L)。',12,True,True,EF),('',4,False,True),
    ('论文实测 (单 CPU, Python, Table 2)：',13,True,False),('',2,False,True),
    ('N=1M: FFD 17s→OBFD 10s (1.7×)',10,False,True,EF),
    ('N=10M: FFD 205s→OBFD 106s (1.93×)',10,False,True,EF),
    ('N=100M: FFD 2311s→OBFD 1066s (2.17×)',10,False,True,EF),
    ('N=1B (10亿): FFD 7.3h→OBFD 3.0h (2.44×)',10,False,True,EF),
    ('N=2B (20亿): FFD 15.3h→OBFD 6.2h (2.48×)',10,False,True,EF),('',3,False,True),
    ('趋势：加速比随 N 单调递增 (1.7×→2.5×)，O(N log L) vs O(N log N) 优势在 N→∞ 时充分体现。3 小时处理 10 亿文档完全满足工业需求。',12,True,True),
])

# ═══════ 11. BF 4/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,11)
Hd(s,'Best-Fit Packing (4/6) — 装箱效率与原理分析')

Img(s,os.path.join(FIGS["bf"],"toy_process_final_trimmed.png"),0.2,1.2,6.3,3.8)
T(s,0.2,5.05,6.3,0.18,'Figure 3: 截断理论分析 — ΔL(m) vs token 位置 m',sz=9,b=False,fn=EF)

Img(s,os.path.join(FIGS["bf"],"truncation_pl_count_2k_only.png"),6.8,1.2,6.2,3.8)
T(s,6.8,5.05,6.2,0.18,'Figure 5: 代码数据截断率统计 (PL)',sz=9,b=False,fn=EF)

M(s,0.75,5.35,11.8,1.85,[
    ('Figure 3 解读：对任意截断概率 p，截断模型 B 的期望损失始终高于完整模型 A。ΔL 随 token 位置 m 单调递增——越靠近文档末尾的 token 受截断损害越大。这解释了截断对代码生成和知识密集型任务影响最严重的原因。',10,False,True),('',2,False,True),
    ('Figure 5 解读：代码数据 L=2048 下截断率与 NL 接近 (~60%)，但代码对截断更敏感——变量定义/使用分离直接导致 Undefined Name 错误。Best-Fit 在程序合成上 +15.0% 绝对提升的根本原因。',10,False,True),
    ('总结：装箱建模→文档完整性从概率约束变为组合优化目标。OBFD→利用 l(c) 有界性实现 O(N log L)。填充率<0.003%→训练效率不变。全维度正向提升无退化→验证了"截断是根因"的假设。',10,True,True),
])

# ═══════ 12. BF 5/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,12)
Hd(s,'Best-Fit Packing (5/6) — 实验结果','7B/13B · 2K/8K · RefinedWeb + The Stack · 22 子任务')

bfr=[['任务类别','提升','代表性数据集 / 说明'],
     ['阅读理解','+4.7%','RACE(+4.3%), SQuADv2(+4.5%), BoolQ(+5.4%), DROP(+4.8%) 等 8 ds'],
     ['自然语言推理','+9.3%','ANLI(+9.4%), HANS, MNLI 共 3 个数据集'],
     ['上下文跟随','+16.8% 最大','22 项中提升最大。直接回应"截断破坏上下文利用"问题'],
     ['程序合成','+15.0% 绝对','HumanEval: 28.3→32.5; MBPP: 35.0→38.1'],
     ['封闭域幻觉','−58.3%','Undefined Name 错误: 3.6%→1.5%'],
     ['退化','0 / 22','全部 22 子任务无统计显著退化']]
Tbl(s,0.4,1.25,12.5,3.0,bfr)

M(s,0.75,4.4,11.8,2.7,[
    ('NTP 训练本质要求"从上下文预测下一个 token"。若 60% 上下文残缺→模型学会"在残缺上下文中预测"。推理时即使给完整上下文，模型仍可能忽略——因为它习惯了。',12,True,True),('',3,False,True),
    ('消灭截断纠正了这一系统性偏差：上下文完整→模型学会基于完整上文做推断→上下文跟随 +16.8%, 幻觉 −58.3%。因果链路清晰。',12,False,True),('',3,False,True),
    ('方法定位：回应缺陷一 (截断)。不改文档顺序/序列长度，仅改装入方式。HF TRL 已集成。局限：不区分语义→ICLM 补足。两者可组合 (先装箱后排序)。',11,False,True),
])

# ═══════ 13. BF 6/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,13)
Hd(s,'Best-Fit Packing (6/6) — 零代价证明与方法定位')

M(s,0.75,1.3,6.2,5.8,[
    ('为什么 Best-Fit 是"零代价"方案？',14,True,False),('',4,False,True),
    ('训练效率维度：填充增加 <0.003%→序列总数几乎不变→每 epoch 训练步数不变→训练总时间不变。GPU 利用率仍为 100%。',11,False,True),('',3,False,True),
    ('预处理维度：O(N log L) 近线性。10 亿文档约 3 小时 (单 CPU)。相对训练本身 (数百 GPU 数周)，预处理时间可忽略不记。',11,False,True),('',3,False,True),
    ('工程集成维度：HuggingFace TRL 已集成 (strategy="ffd")。只需替换数据准备的 packing 函数。不改模型、不改训练循环、不改其他超参。',11,False,True),('',4,False,True),
    ('在三种改进方法中的定位：',14,True,False),('',3,False,True),
    ('改动范围：仅修改"文档→序列"这一步。是最浅层的改进。',11,False,True),
    ('投入产出比：最低投入 (零代价)，最全面收益 (22 子任务无退化)。',11,True,True),
    ('兼容性：与 ICLM 完全正交 (装箱 vs 排序)，可组合使用。',11,False,True),
    ('落地门槛：三种方法中最低。已有成熟开源实现 (HF TRL)。',11,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('从随机拼接到 Best-Fit 的认知转变',14,True,False),('',4,False,True),
    ('随机拼接的核心假设：文档怎么装进序列不重要——GPU 效率优先。',11,False,True),
    ('Best-Fit 用受控实验证伪了这一假设：改变装入方式即可获得全维度正向提升，且无代价。',11,True,True),('',4,False,True),
    ('这一发现将数据拼接从"工程细节"重新定义为"影响模型质量的关键因素"。',12,True,True),('',4,False,True),
    ('Best-Fit 的局限性：',14,True,False),('',3,False,True),
    ('仍不区分文档语义——两个主题无关的文档可能因长度匹配被放入同一序列。跨文档注意力噪声仍然存在 (只是每个文档内部完整了)。',11,False,True),('',3,False,True),
    ('这是 ICLM 的出发点：在保证文档完整性的基础上，通过语义排序让相关文档在一起，进一步消灭跨文档注意力噪声。',12,True,True),
])

# ═══════ 14. ICLM 1/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,14)
Hd(s,'四、In-Context Pretraining (1/6)','Shi, Min et al. ICLR 2024. Meta AI. arXiv:2310.10638')
Img(s,os.path.join(FIGS["ic"],"main.png"),0.2,1.15,8.6,3.9)
T(s,0.2,5.1,8.6,0.18,'Figure 2. ICLM 两步流程：(1)检索构建文档图 (2)图遍历→上下文窗口。',sz=9,b=False)
M(s,9.15,1.2,3.9,4.5,[
    ('Best-Fit 消灭截断，但语义噪声仍存：同一窗口文档随机相遇。ICLM 让语义相关文档出现在同一上下文。仅改排序，不改架构/目标。',10,True,False),('',3,False,True),
    ('Step 1：Contriever 编码→768 维向量。FAISS ANN→top-10 最近邻。',10,False,True),
    ('Step 2：文档图 (节点=文档, 边=互为邻居)→贪心图遍历→排序路径→切为上下文窗口。',10,False,True),('',3,False,True),
    ('核心结果：ICL +8% · 阅读 +15% (最高) · 忠实度 +16% · 检索 +9% · 附带语义去重。效果随数据规模单调递增。',10,True,True),
])

# ═══════ 15. ICLM 2/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,15)
Hd(s,'ICLM (2/6) — 文档嵌入与相似度计算')

Img(s,os.path.join(FIGS["ic"],"intro.png"),0.2,1.2,6.0,2.5)
T(s,0.2,3.75,6.0,0.15,'Figure 1: ICLM vs Standard Pretraining 总览对比',sz=8,b=False,fn=EF)

y=1.2
M(s,6.6,y,6.4,1.5,[
    ('Contriever (Izacard et al., 2022)：无监督对比学习。最后一层隐藏态平均池化：',11,True,False)])
y=F(s,r'E(d)=\frac{1}{|d|}\sum_{t=1}^{|d|}h_t^{(\mathrm{last})}\in\mathbb{R}^{768}',6.6,2.4,6.4,fs=14)
y+=GAP
M(s,6.6,y,6.4,1.0,[('余弦相似度 & 最近邻：',11,True,False)])
y=F(s,r'\mathrm{sim}(d_i,d_j)=\frac{E_i\cdot E_j}{\|E_i\|_2\cdot\|E_j\|_2},\quad N(d_i)=\underset{j\neq i}{\mathrm{top\text{-}k}}\ \mathrm{sim}(d_i,d_j),\ k=10',6.6,y,6.4,fs=11)
y+=GAP
M(s,6.6,y,6.4,1.0,[
    ('FAISS IVF-PQ: 235M docs, 50M/batch, 8 GPU, ~6h。暴力搜索需数周。',10,False,True),
])

M(s,0.2,3.95,6.0,3.2,[
    ('文档图构建：节点=全部文档。边满足双向条件 (至少一方在对方 top-10 中)：',11,True,False)])
F(s,r'(d_i,d_j)\in L\ \Leftrightarrow\ d_j\in N(d_i)\ \text{or}\ d_i\in N(d_j)',0.2,4.5,6.0,fs=12)
M(s,0.2,5.0,6.0,2.0,[
    ('w(d_i,d_j)=sim(d_i,d_j)。双向条件保证连通性。',10,False,True),
    ('图中多个不连通簇，簇内文档主题相关。孤立文档单独成点。',10,False,True),
    ('附带语义去重：sim>0.98→近似重复→自动移除。消融确认去重关键。',10,True,True),
])

# ═══════ 16. ICLM 3/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,16)
Hd(s,'ICLM (3/6) — 贪心图遍历')

M(s,0.75,1.3,6.2,5.8,[
    ('目标：在文档图 G 上找覆盖每个文档恰好一次的哈密顿路径，最大化相邻文档相似度之和：',11,True,False),('',2,False,True)])
F(s,r'\max_P\sum_{t=1}^{N-1}\mathrm{sim}(d_{\pi_t},d_{\pi_{t+1}})',0.75,2.2,6.2,fs=15)
M(s,0.75,2.85,6.2,4.2,[
    ('Maximum TSP (NP-Hard)。N=235M，精确算法不可行。双层贪心近似：',11,True,True),('',3,False,True),
    ('外层 — 选链段起点：每次从未访问节点中选度数最小的。度数低=相似文档少。优先处理避免最后孤立。最小度数优先启发式 (Min-Degree-First)。',11,True,True),('',3,False,True),
    ('内层 — 扩展链段：从当前节点出发，在未访问邻居中选余弦相似度最高的作为下一步。重复到当前连通分量全部访问 (死胡同)→回外层选新起点。',11,True,True),('',3,False,True),
    ('复杂度：每条边在检查时最多访问一次。|L|≈N·k/2≈1.2B 边→O(N·k)=O(N)。20 CPU×12h (235M 文档)。',11,True,True),('',3,False,True),
    ('路径→上下文窗口：沿路径累积 token 到 L 切窗口。batch 内取路径不同区段保多样性。相邻文档语义相关 → 跨文档注意力从噪声变信号。',11,True,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('为什么从最小度数节点开始？',13,True,False),('',3,False,True),
    ('图的度数分布不均匀。热门文档度数极高，冷门文档度数极低。若最后留冷门文档，其邻居可能已全部访问→成为孤立点→只能零权边跳转→路径高权边少→语义连贯性差。',10,False,True),('',3,False,True),
    ('最小度数优先→先处理冷门→冷门节点的邻居大概率也未访问→可在它们之间形成短但连贯的链段。热门节点度数高，晚处理也能找到出路。',10,True,True),('',3,False,True),
    ('论文验证：最小度数起点 vs 随机起点，前者路径的平均余弦相似度显著更高 (0.35 vs 0.27，随机排列仅 0.02)。',10,False,True),('',4,False,True),
    ('计算开销：检索 32GPU×6h + 图遍历 20CPU×12h + 训练 128A100×9d。预处理占训练总成本 <2%。',10,True,True),('',3,False,True),
    ('方法定位：回应缺陷二 (噪声)。可与 Best-Fit 组合。需跨文档推理 (RAG/ICL/多跳) 受益最大。',10,True,True),
])

# ═══════ 17. ICLM 4/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,17)
Hd(s,'ICLM (4/6) — k 值分析与语义去重')

M(s,0.75,1.3,6.2,5.8,[
    ('k (top-k 最近邻) 对图结构的影响',14,True,False),('',3,False,True),
    ('k 控制文档图的连接密度——这是 ICLM 最关键的超参数。',12,True,True),('',3,False,True),
    ('k=1：每文档仅连最相似 1 个邻居。图极度稀疏，大量孤立节点和小型连通分量。路径频繁断裂→大量零权边→语义连贯性最差。',10,False,True),('',3,False,True),
    ('k=5：每文档连 5 个邻居 (+被连)。图的中等密度。多数文档形成中等大小连通分量。路径断裂减少但仍有较多零权边。',10,False,True),('',3,False,True),
    ('k=10 (论文最优)：每文档连 10 个邻居。图在连接性和多样性之间取得平衡。各语义主题形成稳定的连通分量。路径中零权边最少，语义连贯性最强。',10,True,True),('',3,False,True),
    ('k=20：图过于稠密。高流行度文档被过度连接 (度数≈k²)，图遍历路径被少数热门文档主导，丧失语义多样性。训练效率下降 (路径倾向于停留在少量高密度区域)。',10,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('语义去重 — ICLM 的附带收益',14,True,False),('',3,False,True),
    ('在 FAISS 检索阶段，论文发现近似重复文档 (同一文档的多个镜像/轻微编辑版) 产生极高的余弦相似度 (sim>0.98)。',10,False,True),('',3,False,True),
    ('这些近似重复对在文档图中自然形成密集子图——它们之间的相似度极高，且与外部文档的相似度分布也高度一致。ICLM 可以在构建图的过程中直接识别这些子图并去重。',10,False,True),('',3,False,True),
    ('消融实验确认：不做语义去重将导致 ICLM 效果显著退化。近重复文档扭曲了图的度数分布 (这些文档的度数异常高)，图遍历路径被过度引导向这些重复区域，降低了路径的语义多样性。',10,True,True),('',3,False,True),
    ('这证明了：语义去重不是可选的附加步骤，而是 ICLM pipeline 的必要组成部分。论文在最终方案中包含了这一步骤。',11,True,True),('',3,False,True),
    ('消融总结：k=10 最优；BM25 替代 Contriever 效果下降 (短文档和代码尤甚)；不去重显著退化。三项消融共同确认了 ICLM 各组件设计的合理性。',11,False,True),
])

# ═══════ 18. ICLM 5/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,18)
Hd(s,'ICLM (5/6) — 实验结果','300B tokens CommonCrawl · 7B · L=8192')

icr=[['任务类别','提升','数据集 / 说明'],
     ['In-Context Learning','+8%','7 分类数据集 (32-shot)。全部优于 Standard。与 kNN 持平且无重复'],
     ['阅读理解','+15% 最高','SQuADv2, RACE, BoolQ, DROP, HotpotQA 等 6 ds (2-shot)。四种方法最高'],
     ['忠实度','+16%','NQ-Swap, MemoTrap。更倾向于遵循给定上下文而非依赖预训练记忆'],
     ['检索增强','+9%','Wikipedia 外部知识的开放域 QA'],
     ['长上下文推理','+5%','长上下文基准测试'],
     ['困惑度','全部更低','Wikipedia/Arxiv/Books, 0.3B–7B 全部规模优于 Standard']]
Tbl(s,0.3,1.25,12.7,2.8,icr)

Img(s,os.path.join(FIGS["ic"],"ppl.png"),0.3,4.2,5.3,1.5)
T(s,0.3,4.08,5.3,0.12,'Figure 3: 各模型规模 ICLM PPL 均低于 Standard',sz=7,b=False,fn=EF)

Img(s,os.path.join(FIGS["ic"],"evolution.png"),5.9,4.2,7.1,1.5)
T(s,5.9,4.08,7.1,0.12,'Figure 4: ICLM 效果随训练过程的演化',sz=7,b=False,fn=EF)

M(s,0.3,5.8,12.7,1.3,[
    ('扩展性：数据 100M→10B tokens, ICLM 优势 +4%→+11%，单调递增无饱和。更大规模下优势可能继续扩大。',10,False,True),
    ('阅读理解提升最大 (+15%) 的原因：ICLM 预训练输入 = 多篇相关文档 → 天然"综合多文档"格式 → 与 2-shot ICL 评估分布对齐。这是预训练与推理分布对齐 (Distribution Alignment) 的直接体现。',11,True,True),
])

# ═══════ 19. ICLM 6/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,19)
Hd(s,'ICLM (6/6) — 方法定位与组合使用')

M(s,0.75,1.3,6.2,5.8,[
    ('ICLM 的独特位置',14,True,False),('',3,False,True),
    ('四种方法中唯一从"语义层面"解决拼接质量的方法。与 Best-Fit (长度层面) 完全正交——前者决定空间安排 (哪些 chunk 放一起)，后者决定时间顺序 (文档排列的先后)。两者在同一训练管线中共存。',11,False,True),('',4,False,True),
    ('ICLM 与 In-Context Learning 的深层关系',14,True,False),('',3,False,True),
    ('ICL (推理) 要求模型利用上下文中的几个示例快速适应新任务。这个能力的前提：模型在预训练中就学会了"阅读并利用上下文中的相关信息"。',11,False,True),('',3,False,True),
    ('随机拼接：上下文内文档随机→没有练习跨文档推理→推理时做 ICL 需要的能力是预训练中未训练过的。',11,False,True),
    ('ICLM：上下文内文档相关→频繁练习跨文档推理→推理时做 ICL 的能力正是预训练中被大量训练的。',11,True,True),('',3,False,True),
    ('同样逻辑解释忠实度 +16%：当模型习惯上下文中的文档是可靠信息源时，更倾向于依赖给定上下文而非预训练记忆。这对 RAG 场景至关重要 (需要引用检索到的文档)。',11,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('工业可行性分析',14,True,False),('',3,False,True),
    ('优点：仅改变排序，不改模型架构/训练目标。可直接集成到现有预训练管线。附带语义去重 (无需额外计算)。',11,False,True),('',3,False,True),
    ('门槛：(1) 需要 GPU 集群做嵌入和检索 (32 GPU)。虽 <2% 训练总成本但增加基础设施需求。',11,False,True),
    ('(2) 依赖 Contriever 嵌入质量。对非英语语料可能需要重新训练嵌入模型 (论文明确的局限性)。',11,False,True),
    ('(3) 预处理离线完成，多 epoch 训练时数据顺序不变。如需不同 epoch 用不同排序→需重新检索和遍历。',11,False,True),('',4,False,True),
    ('组合使用建议',14,True,False),('',3,False,True),
    ('Best-Fit + ICLM：Phase 1 Best-Fit 装箱保证完整性→Phase 2 ICLM 排序增强语义连贯性。两者完全正交，可叠加。',11,True,True),
    ('DD + ICLM：桶内做语义排序 (相关主题文档片段在训练时间上靠近)。桶间由 VSL 课程调度决定采样顺序。',11,False,True),
])

# ═══════ 20. DD 1/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,20)
Hd(s,'五、Dataset Decomposition (1/6)','Pouransari, Li et al. NeurIPS 2024. Apple + Anthropic. arXiv:2405.13226')
Img(s,os.path.join(FIGS["dd"],"x2.png"),0.2,1.15,7.5,3.9)
T(s,0.2,5.1,7.5,0.18,'二进制分解示意：文档按 l=Σ2^{i_k} 切分，各入对应桶 (D₈~D₁₃)。',sz=9,b=False)
M(s,8.0,1.2,5.0,3.5,[
    ('前三种方法回答"如何把文档拼成固定 L"。DD 问：L 为什么要固定？',10,True,False),('',2,False,True),
    ('二进制分解 — 文档按长度 l 的二进制展开切分：',10,False,True)])
F(s,r'l=\sum_{k=1}^{K}2^{i_k},\quad i_1>i_2>\cdots>i_K\geq 0',8.0,2.8,5.0,fs=15)
M(s,8.0,3.5,5.0,3.6,[
    ('长度 2^{i_k}→桶 D_{i_k}。i∈{8..13}={256,512,1024,2048,4096,8192}。',10,True,True),
    ('例: 2200=2048+128+16+8→D₁₁,D₇,D₄,D₃',10,False,True,EF),
    ('例: 5600=4096+1024+256+128+64+32',10,False,True,EF),
    ('桶内序列同长→天然对齐。每序列=单文档片段→零跨文档注意力。',10,True,True),
    ('核心结果：6×加速 · 4×数据效率 · 训8K≈训2K · Loss 3.18。',11,True,True),
])

# ═══════ 21. DD 2/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,21)
Hd(s,'DD (2/6) — 二进制分解原理与 VSL 训练')
y=1.3
M(s,0.75,y,6.2,1.5,[('选择 2 的幂次的四个理由：',13,True,False),
    ('(1) GPU 矩阵乘法对齐：cuBLAS tile=128/256，2 的幂次序列避免最后 tile 部分填充。',10,True,True),
    ('(2) 二进制最优性：任意整数唯一二进制表示，1-bit 数 (分段数) 在所有 base-B 展开中最少。',10,True,True),
    ('(3) FlashAttention 兼容：块内的 token 完整分割，避免跨块通信。',10,True,True),
    ('(4) 桶数可控：仅 6 个桶 (2⁸~2¹³)，每桶序列充足，batch 采样稳定。',10,True,True),])
y=GAP+3.0
M(s,0.75,y,6.2,1.5,[('VSL 训练 — 每步 token 预算 b (常数)。采样桶 D_i：',11,True,False)])
y=F(s,r'\mathrm{BS}_i=b\,/\,2^i,\qquad \mathrm{BS}_i\times 2^i=b',0.75,y+0.2,6.2,fs=15)
y+=GAP
M(s,0.75,y,6.2,1.0,[('每步 FLOPs 推导：自注意力二次于序列长度。单序列 FLOPs∝(2^i)²：',10,True,False)])
y=F(s,r'\mathrm{FLOPs}_i\propto\mathrm{BS}_i\times(2^i)^2=\frac{b}{2^i}\times 2^{2i}=b\times 2^i',0.75,y,6.2,fs=14)
y+=GAP
M(s,0.75,y,6.2,1.5,[
    ('每步 FLOPs 与序列长度成正比——而非二次！(batch size 与长度互为倒数，消去了一个 2^i)',10,True,True),
    ('D₈(256): b×256 | D₁₁(2048): b×2048 | D₁₃(8192): b×8192 (32× 差距)',10,False,True,EF),
])

M(s,7.2,1.3,5.7,5.8,[
    ('传统固定 L=8192：每步恒 b×8192。DD 加权平均远低于此。',11,True,True),('',3,False,True),
    ('Grow-P2 课程策略：采样权重由 α 控制',13,True,False)])
F(s,r'w_i\propto\frac{1}{2^{\,\alpha\,(i-i_{\min})}}',7.2,2.65,5.7,fs=16)
M(s,7.2,3.4,5.7,3.7,[
    ('α 随时间递减。早期 α 大→w 极度偏向短桶 (廉价基础训练)。',11,False,True),
    ('后期 α 小→长桶频率增加 (长上下文精训)。',11,False,True),
    ('论文分布: 早期 ~70% D₈–D₁₀, ~20% D₁₁–D₁₂, ~10% D₁₃。中后期拉平。',10,False,True,EF),('',3,False,True),
    ('加权平均 FLOPs ≈ b×(0.70·256+0.20·2048+0.10·8192) ≈ b×1408',11,True,True,EF),
    ('加速比: 8192/1408≈5.8≈6×',14,True,True),
    ('传统 2K: b×2048。DD 加权 b×1408 < b×2048',11,True,True),
    ('→ DD 训 8K 总成本 < 传统训 2K 的成本！',13,True,True),('',3,False,True),
    ('消融 (论文 Table 2): DD+Grow-P2 Loss 3.18 > DD+Linear 3.22 > DD+Random 3.35 > 基线 3.55。课程贡献约一半增益 (~46%)。',10,True,True),
])

# ═══════ 22. DD 3/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,22)
Hd(s,'DD (3/6) — 1000 文档示例与 VSL 训练分析')

M(s,0.75,1.3,6.2,5.8,[
    ('1000 文档完整示例 (μ=2000, σ=1500, Σ≈2,000,000 tokens)：',13,True,False),('',3,False,True),
    ('二进制分解后桶分布：',12,True,True),('',2,False,True),
    ('D₁₃(8192):~45 条 → ~368K (18.4%) | D₁₂(4096):~120 → ~492K (24.6%)',10,False,True,EF),
    ('D₁₁(2048):~280 → ~573K (28.7%) ← token 最多',10,False,True,EF),
    ('D₁₀(1024):~350 → ~358K (17.9%) | D₉(512):~300 → ~154K (7.7%)',10,False,True,EF),
    ('D₈(256):~200 → ~51K (2.6%)',10,False,True,EF),('',2,False,True),
    ('共~1295 序列, ~1,996,800 tokens。零丢弃, 零 padding。',11,True,True),('',4,False,True),
    ('VSL 训练 (b=2,000,000) 各桶 step 特征：',12,True,True),
    ('D₁₃: BS≈244, FLOPs∝b×8192 (与传统相同)',10,False,True,EF),
    ('D₁₁: BS≈976, FLOPs∝b×2048',10,False,True,EF),
    ('D₈: BS≈7812, FLOPs∝b×256 (极低廉，大 batch→稳定梯度)',10,False,True,EF),
])

M(s,7.2,1.3,5.7,5.8,[
    ('DD 与前三种方法在文档处理上的根本区别：',13,True,False),('',3,False,True),
    ('随机拼接：文档被位置随机截断→不完整片段+噪声。',10,False,True),
    ('Best-Fit：文档按长度 (l>L) 确定切分。l≤L 完整。序列可含多文档。',10,False,True),
    ('ICLM：Best-Fit 基础上做语义排序。文档处理方式不变。',10,False,True),
    ('DD：主动切分所有文档为 2 的幂次片段。切分是精心设计的 (GPU 友好的长度, 来自同一文档连续区域)。不同文档片段永不混在同一序列中。',10,True,True),('',4,False,True),
    ('DD 的序列数自然多于原始文档数 (1000→1295)。但在 VSL 训练中，大多数步来自短桶 (D₈-D₁₀)，这些步 batch size 极大、FLOPs 极低。加权后总训练时间反而更短。这是 DD 6× 加速的数学基础。',10,True,True),('',4,False,True),
    ('方法定位：打破从 GPT-1 起"序列长度必须固定"的根本范式。与 BF/ICLM 兼容——桶内可再做装箱优化和语义排序。需修改训练循环 (变长 batch + 课程调度)。课程不可省略。已开源 (github.com/apple/ml-dataset-decomposition)。',10,True,True),
])

# ═══════ 23. DD 4/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,23)
Hd(s,'DD (4/6) — Grow-P2 课程深入分析')

M(s,0.75,1.3,6.2,5.8,[
    ('Grow-P2 为什么有效？—— 课程学习的理论直觉',14,True,False),('',3,False,True),
    ('机器学习中的课程学习 (Curriculum Learning, Bengio et al., 2009) 核心思想：先学简单样本，逐步过渡到复杂样本。DD 将这一思想引入了序列长度维度。',11,False,True),('',3,False,True),
    ('短序列步 (<1024 tokens)：覆盖词频统计、语法结构、短程依赖。这些"基础能力"对序列长度不敏感——512 token 足以覆盖绝大多数语法现象。大量廉价训练步骤提供了稳定的基础语言能力。',10,True,True),('',3,False,True),
    ('中等序列步 (2048–4096)：增加段落级上下文。模型开始学习跨句依赖、局部篇章结构。此时基础已经打牢，增加长上下文不会导致训练不稳定。',10,False,True),('',3,False,True),
    ('长序列步 (8192)：集中学习长程依赖、跨段落推理。这些步骤数量少但效果显著——因为基础能力已在廉价训练中打好。类比：先学好基础课再上专业课，比一开始就上专业课更高效。',10,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('Grow-P2 vs 其他课程策略的消融分析',14,True,False),('',3,False,True),
    ('Baseline (Concat-and-Chunk, 固定 L=8192)：Loss=3.55。无 VSL，无课程。所有训练步都使用最长的序列长度，每步成本最高。',10,False,True,EF),('',2,False,True),
    ('DD + Random：各桶按 token 量均匀随机采样。Loss=3.35。VSL 训练本身降低 Loss 0.20——长度可变训练的泛化能力优于固定长度训练。但无课程调度时，长序列步在早期就占用大量算力，效率低下。',10,False,True,EF),('',2,False,True),
    ('DD + Linear：线性调整各桶权重。Loss=3.22。有调度优于随机调度，但线性不够精细——早期长桶权重仍然偏高。',10,False,True,EF),('',2,False,True),
    ('DD + Grow-P2：2 幂次递减权重。Loss=3.18 ✓ (最优)。权重曲线形状恰好匹配训练的自然需求——早期极度偏向短桶 (指数级倾斜)，后期逐步拉平。',10,True,True,EF),('',3,False,True),
    ('课程贡献分析：Grow-P2 在 Random 基础上再降 Loss 0.17。贡献了约 46% 的总增益 (0.17/0.37)。课程不是锦上添花——没有课程，DD 的一半增益将丢失。这是 DD 论文最重要的消融发现。',12,True,True),
])

# ═══════ 24. DD 5/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,24)
Hd(s,'DD (5/6) — 实验结果','RefinedWeb (525B) + DCLM · 160M–1B · L_max=8192')

ddr=[['指标','结果','说明'],
     ['训练加速','6×','达到相同 Loss 仅需 1/6 训练 token (time-to-accuracy)'],
     ['数据效率','4×+','优势单调递增: 10B +3%→200B +15%→1.1T +22%。未饱和'],
     ['长上下文成本','训 8K≈训 2K','RULER: DD(8K) 73.0 vs 基线(2K) 55.1 vs 基线(8K) 62.3'],
     ['消融 (Grow-P2)','Loss 3.18 最优','DD+Grow-P2 3.18 > DD+Linear 3.22 > DD+Random 3.35 > 基线 3.55'],
     ['课程贡献','~46% 总增益','Grow-P2 在 Random 基础上再降 0.17 Loss。课程为必须组件'],
     ['开源','GitHub','github.com/apple/ml-dataset-decomposition']]
Tbl(s,0.3,1.25,12.7,2.6,ddr)

Img(s,os.path.join(FIGS["dd"],"x4.png"),0.3,4.0,5.5,1.6)
T(s,0.3,3.88,5.5,0.12,'训练效率曲线: DD vs Concat-and-Chunk 基线',sz=7,b=False,fn=EF)

Img(s,os.path.join(FIGS["dd"],"x7.png"),6.1,4.0,6.9,1.6)
T(s,6.1,3.88,6.9,0.12,'消融: Grow-P2 vs Linear vs Random 课程策略',sz=7,b=False,fn=EF)

M(s,0.3,5.7,12.7,1.4,[
    ('方法定位：DD 是唯一打破"固定 L"范式的方法。深度创新但落地门槛高——需修改训练循环。与 BF/ICLM 兼容 (桶内可装箱+排序)。适用场景：预算有限需高效训长上下文；或追求训练效率最大化。',10,False,True),
    ('局限：二进制分解增加预处理。若完全不训长序列→无长上下文能力→后期必须包含足够长序列比例。对非 2 幂次 L 需 padding。验证规模 1B，百亿+规模尚未验证。',10,False,True),
])

# ═══════ 25. DD 6/6 ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,25)
Hd(s,'DD (6/6) — 范式意义与组合使用')

M(s,0.75,1.3,6.2,5.8,[
    ('DD 的范式意义',14,True,False),('',3,False,True),
    ('从 2018 年 GPT-1 起，LLM 预训练的一个基本前提是"序列长度必须固定"。这是 GPU 硬件约束的直接推论。所有人都在这个前提下做优化——如何更好地把文档装入固定长度的容器。',11,False,True),('',3,False,True),
    ('DD 证明了一个被长期忽视的事实：序列长度可以不固定。通过 VSL 训练 (序列长度与 batch size 互为倒数)，每步总 token 数恒为 b，GPU 效率不受损失。固定长度范式是历史惯性的产物，不是计算的必然。',11,True,True),('',3,False,True),
    ('这一洞察的影响超出了数据拼接本身。它为训练长上下文模型提供了一种全新的经济可行的路径：不需要支付 O(L²) 的全额代价，通过课程调度可以将长上下文训练成本降至传统短上下文水平。',11,False,True),('',4,False,True),
    ('DD 的局限与未来方向',14,True,False),('',3,False,True),
    ('当前验证在 1B 规模。百亿+参数规模的验证尚未进行——DD 在超大模型上的行为有待验证。二进制分解产生的碎片 (一篇文档拆为多段) 是否在更大规模上产生新的问题 (如长文档的跨片段一致性) 有待研究。',10,False,True),
])

M(s,7.2,1.3,5.7,5.8,[
    ('三种方法的组合使用路线图',14,True,False),('',3,False,True),
    ('基础层 — Best-Fit Packing：在所有场景优先采用。零代价消灭不必要截断。已集成 HF TRL，开箱即用。提供文档完整性的基本保证。',11,True,True),('',3,False,True),
    ('增强层 — ICLM：在需要跨文档推理的下游任务上叠加。RAG、Few-shot ICL、多跳推理受益最大。需额外 GPU 设施 (检索)。与 Best-Fit 完全正交——Phase 1 装箱保证完整性→Phase 2 排序增强连贯性。',11,True,True),('',3,False,True),
    ('范式层 — DD：面向未来高效长上下文训练。与 Best-Fit 和 ICLM 均兼容——在桶内对序列做装箱优化和语义排序。需修改训练循环但总训练成本大幅降低。目前最前沿但落地门槛最高。',11,True,True),('',3,False,True),
    ('选择建议：从 Best-Fit 开始 (最低投入)→按需叠加 ICLM (需跨文档推理时)→关注 DD 进展 (高效长上下文是大趋势，等待百亿+规模验证后大规模推广)。',12,True,True),
])

# ═══════ 26-28. COMPARE ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,26)
Hd(s,'六、四种方法对比 — 方法论 (1/3)')
cd1=[['维度','随机拼接','Best-Fit Packing','ICLM','Dataset Decomp.'],
     ['操作','打乱→EOS→按L切分','l>L切chunk, BFD装箱','Contriever→FAISS→图遍历','二进制分解→入桶→VSL'],
     ['优化目标','GPU利用率(零padding)','文档完整性(最小截断)','语义连贯性(噪声→信号)','训练成本(6×加速)'],
     ['截断率(L=2048)','~60%','<5%(仅l>L必要截断)','<1%','100%主动分解(非截断)'],
     ['跨文档注意力','大量噪声(~15-30%)','存在(chunk间)','存在(相关,可能有用)','零(单文档片段)'],
     ['序列长度','固定L','大多为L','大多为L','{256,...,8192},动态'],
     ['额外开销','无','极小 O(N log L)','中等(检索+图遍历)','小(分解+课程)'],
     ['工程成熟度','所有框架原生','HF TRL已集成','需GPU集群(检索)','Apple已开源(OpenLM)'],
     ['发表','GPT-3,LLaMA,PaLM等','Ding et al. ICML 2024 Amazon','Shi et al. ICLR 2024 Meta','Pouransari et al. NeurIPS 2024 Apple']]
R,_=len(cd1),len(cd1[0])
tb=s.shapes.add_table(R,_,Inches(0.2),Inches(1.25),Inches(12.9),Inches(5.75)).table
for r in range(R):
    for c in range(_):
        cl=tb.cell(r,c);cl.text=cd1[r][c]
        for p in cl.text_frame.paragraphs:
            p.font.size=Pt(10);p.font.name=CF
            if r==0:p.font.bold=True;p.font.color.rgb=WH;p.alignment=PP_ALIGN.CENTER
            elif c==0:p.font.bold=True;p.font.color.rgb=BK
            else:p.font.color.rgb=BK
        cl.fill.solid();cl.fill.fore_color.rgb=BL if r==0 else (LG if r%2==0 else WH)

s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,27)
Hd(s,'六、四种方法对比 — 性能与递进关系 (2/3)')
cd2=[['指标','随机拼接','Best-Fit Packing','ICLM','Dataset Decomp.'],
     ['训练效率','100%','≈100%(<0.003%)','≈100%','加权约6×更快(vs L=8192)'],
     ['阅读理解','基线','+4.7%(8 ds)','+15%(6 ds)←最高','Loss 3.18 vs 3.55'],
     ['上下文/幻觉','基线','+16.8% / −58.3%','忠实度+16% / 检索+9%','RULER 73.0(vs 55.1)'],
     ['数据效率','1×','≈1×','≈1×','4×+(单调↑,未饱和)'],
     ['长上下文成本','O(L²);L×4→×16','同左','同左','训8K≈传统2K'],
     ['递进关系','起点:发现截断','回应缺陷一(截断)\n装箱消灭不必要截断\n零代价,可组合','回应缺陷二(噪声)\n语义排序→信号\n与BF正交','回应范式(固定L)\n不再要求L固定\n与前两者兼容'],
     ['推荐场景','快速原型,管线验证','正式预训练(百亿+)\n最成熟改进方案','跨文档推理:RAG/ICL\n多跳推理,检索增强','高效长上下文训练\n预算受限,效率优先']]
R2,_=len(cd2),len(cd2[0])
tb2=s.shapes.add_table(R2,_,Inches(0.2),Inches(1.25),Inches(12.9),Inches(5.75)).table
for r in range(R2):
    for c in range(_):
        cl=tb2.cell(r,c);cl.text=cd2[r][c]
        for p in cl.text_frame.paragraphs:
            p.font.size=Pt(10);p.font.name=CF
            if r==0:p.font.bold=True;p.font.color.rgb=WH;p.alignment=PP_ALIGN.CENTER
            elif c==0:p.font.bold=True;p.font.color.rgb=BK
            else:p.font.color.rgb=BK
        cl.fill.solid();cl.fill.fore_color.rgb=BL if r==0 else (LG if r%2==0 else WH)

s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,28)
Hd(s,'六、总结与演进 (3/3)')
tl=[('2018–2023','GPT-3 (2020,175B,L=2048)、LLaMA (2023,7B–65B)、PaLM (2022,540B) 等统一使用随机拼接。五年间被视为已解决的工程细节。Ding et al. (2024) 统计截断率约 60%。'),
    ('2024.01 ICLR','Shi et al. (Meta) — ICLM。Contriever+FAISS+贪心图遍历。仅改排序不改架构。阅读+15%,忠实度+16%,附带语义去重。'),
    ('2024.04 ICML','Ding et al. (Amazon) — Best-Fit Packing。装箱+BFD+线段树 O(N log L)。零代价截断 60%→~5%，幻觉−58.3%。HF TRL 已集成。'),
    ('2024.05 NeurIPS','Pouransari et al. (Apple) — Dataset Decomposition。二进制分解+VSL+Grow-P2 打破固定 L 范式。6× 加速，训8K≈传统2K。开源。')]
for i,(t,desc) in enumerate(tl):
    y=1.3+i*1.5;T(s,0.85,y,3.3,0.35,t,sz=16,b=True,fn=EF)
    T(s,0.85,y+0.45,11.5,0.9,desc,sz=11,b=False)
    if i<3:Sep(s,y+1.42)

# ═══════ 29. REFS ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s);Bb(s);Pn(s,29)
Hd(s,'参考文献')
refs=['[1]  Ding, Wang et al. "Fewer Truncations Improve Language Modeling." ICML 2024. arXiv:2404.10830',
    '[2]  Shi, Min et al. "In-Context Pretraining." ICLR 2024. arXiv:2310.10638',
    '[3]  Pouransari, Li et al. "Dataset Decomposition." NeurIPS 2024. arXiv:2405.13226',
    '[4]  Brown et al. "Language Models are Few-Shot Learners." NeurIPS 2020.',
    '[5]  Touvron et al. "LLaMA: Open and Efficient Foundation Language Models." 2023.',
    '[6]  Chowdhery et al. "PaLM: Scaling Language Modeling with Pathways." 2022.',
    '[7]  Zhao et al. "Analysing The Impact of Sequence Composition on LM Pre-Training." ACL 2024.',
    '[8]  Izacard et al. "Unsupervised Dense Information Retrieval with Contrastive Learning." TMLR 2022.',
    '[9]  Johnson et al. "Billion-scale similarity search with GPUs." IEEE Big Data 2019.',
    '[10] Zhuang et al. "LongPack." ICLR 2025.',
    '[11] Yin et al. "Seamless Data Packing." ACL 2025 Findings.',
    '[12] Dai et al. "DeepSeek-V2." 2024. arXiv:2405.04434']
for i,ref in enumerate(refs):T(s,0.55,1.45+i*0.44,11.8,0.4,ref,sz=10,b=False)

# ═══════ 30. THANK YOU ═══════
s=prs.slides.add_slide(prs.slide_layouts[6]);Bg(s)
Ln(s,0,0,13.333,7.5,fc=BL);Ln(s,-1.5,3.8,6.5,6.5,fc=DL)
T(s,0,2.2,13.333,1.0,'感谢聆听',sz=52,c=WH,b=True,al=PP_ALIGN.CENTER)
T(s,0,3.2,13.333,0.5,'Thank You',sz=24,c=RGBColor(0xBB,0xCC,0xF0),al=PP_ALIGN.CENTER,fn=EF)
Ln(s,5.8,4.0,1.8,0.025,fc=WH)
T(s,0,4.3,13.333,0.35,'大模型预训练数据拼接策略',sz=16,c=RGBColor(0x99,0xB0,0xE0),al=PP_ALIGN.CENTER)
T(s,0,5.0,13.333,0.3,'2026 年 6 月',sz=13,c=RGBColor(0x99,0xB0,0xE0),al=PP_ALIGN.CENTER)

out=os.path.join(BASE,'数据拼接策略_终稿.pptx')
prs.save(out)
tmp=os.path.join(BASE,'_tf.png')
if os.path.exists(tmp):os.remove(tmp)
print(f'OK: {len(prs.slides)} slides')
